# ingest_file.py
import os
import sys
from pinecone import Pinecone
from openai import OpenAI
from dotenv import load_dotenv

# File-specific libraries
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

# Load environment variables from .env file
load_dotenv()

# Log environment variables for debugging (avoid printing sensitive keys in production)
print("OPENAI_API_KEY:", os.getenv("OPENAI_API_KEY"))
print("PINECONE_API_KEY:", os.getenv("PINECONE_API_KEY"))
print("PINECONE_INDEX_NAME:", os.getenv("PINECONE_INDEX_NAME"))
print("PINECONE_NAMESPACE:", os.getenv("PINECONE_NAMESPACE", "default"))

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
PINECONE_NAMESPACE = os.getenv("PINECONE_NAMESPACE", "default")

client = OpenAI(api_key=OPENAI_API_KEY)
pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(PINECONE_INDEX_NAME)

def extract_text_from_pdf(file_path):
    print(f"Extracting text from PDF: {file_path}")
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        t = page.extract_text()
        if t:
            text += t + "\n"
    return text

def extract_text_from_pptx(file_path):
    print(f"Extracting text from PPTX: {file_path}")
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path):
    print(f"Extracting text from DOCX: {file_path}")
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def chunk_text(full_text):
    chunks = [chunk.strip() for chunk in full_text.split("\n\n") if chunk.strip()]
    print(f"Chunked text into {len(chunks)} chunks")
    return chunks

def generate_vectors(chunks):
    vectors = []
    for i, chunk in enumerate(chunks):
        try:
            print(f"Generating embedding for chunk {i} (length: {len(chunk)})")
            response = client.embeddings.create(
                model="text-embedding-ada-002",
                input=chunk
            )
            embedding = response.data[0].embedding  # Using attribute access
            vectors.append({
                "id": f"chunk-{i}",
                "values": embedding,
                "metadata": {"text": chunk}
            })
        except Exception as e:
            print(f"Error processing chunk {i}: {e}")
    print(f"Generated {len(vectors)} vectors")
    return vectors

def upsert_vectors(vectors):
    if vectors:
        upsert_response = index.upsert(vectors=vectors, namespace=PINECONE_NAMESPACE)
        print("Upserted vectors response:", upsert_response)
    else:
        print("No vectors to upsert.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ingest_file.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    print("Ingesting file:", file_path)
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        full_text = extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        full_text = extract_text_from_pptx(file_path)
    elif ext == ".docx":
        full_text = extract_text_from_docx(file_path)
    else:
        raise ValueError("Unsupported file type: " + ext)

    print("Extracted text length:", len(full_text))
    chunks = chunk_text(full_text)
    vectors = generate_vectors(chunks)
    upsert_vectors(vectors)
